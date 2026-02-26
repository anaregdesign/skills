#!/usr/bin/env python3
"""Append one slide to an existing PowerPoint file from a JSON slide spec."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

from build_pptx import (
    add_notes,
    add_table,
    add_visual,
    list_master_layouts,
    load_json,
    looks_like_json_object,
    normalize_slide_spec,
    pick_layout,
    set_bullets,
    set_subtitle,
    set_title,
)

JAPANESE_CHAR_RE = re.compile(r"[ぁ-ゖァ-ヺー一-龯々〆〤]")
ASCII_LETTER_RE = re.compile(r"[A-Za-z]")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--deck", required=True, help="Target PPTX path to append to")
    parser.add_argument(
        "--kind",
        choices=["title", "content"],
        default="content",
        help="Slide spec kind: title or content (default: content)",
    )
    parser.add_argument(
        "--spec",
        required=True,
        help="Slide spec JSON file path",
    )
    parser.add_argument(
        "--plan-base",
        help="Base directory for resolving relative paths",
    )
    parser.add_argument(
        "--fallback-layout",
        type=int,
        default=1,
        help="Deprecated (ignored). Layout is always auto-detected from slide masters.",
    )
    parser.add_argument(
        "--strict-images",
        action="store_true",
        help="Fail when a referenced visual image file is missing",
    )
    parser.add_argument(
        "--no-prefer-master-layouts",
        action="store_true",
        help="Deprecated (ignored). Master layout auto-detection is always enabled.",
    )
    parser.add_argument(
        "--language",
        help="Expected slide language: ja or en",
    )
    return parser.parse_args()


def normalize_language(raw: Any) -> str | None:
    if raw is None:
        return None
    value = str(raw).strip().lower()
    if value in {"ja", "jp", "japanese"}:
        return "ja"
    if value in {"en", "english"}:
        return "en"
    return None


def contains_japanese(text: str) -> bool:
    return bool(JAPANESE_CHAR_RE.search(text))


def contains_ascii_letters(text: str) -> bool:
    return bool(ASCII_LETTER_RE.search(text))


def collect_slide_texts(raw_spec: dict[str, Any]) -> list[str]:
    values: list[str] = []
    for key in ("title", "subtitle"):
        raw = raw_spec.get(key)
        if raw is not None:
            text = str(raw).strip()
            if text:
                values.append(text)

    bullets = raw_spec.get("bullets", [])
    if isinstance(bullets, list):
        for item in bullets:
            text = str(item).strip()
            if text:
                values.append(text)

    visual = raw_spec.get("visual")
    if isinstance(visual, dict):
        caption = str(visual.get("caption", "")).strip()
        if caption:
            values.append(caption)

    table = raw_spec.get("table")
    if isinstance(table, dict):
        headers = table.get("headers", [])
        rows = table.get("rows", [])
        if isinstance(headers, list):
            for header in headers:
                text = str(header).strip()
                if text:
                    values.append(text)
        if isinstance(rows, list):
            for row in rows:
                if not isinstance(row, list):
                    continue
                for cell in row:
                    text = str(cell).strip()
                    if text:
                        values.append(text)
    return values


def validate_slide_language(raw_spec: dict[str, Any], expected_language: str | None) -> None:
    if expected_language is None:
        return
    texts = collect_slide_texts(raw_spec)
    if not texts:
        return

    ja_count = sum(1 for text in texts if contains_japanese(text))
    en_count = sum(
        1
        for text in texts
        if contains_ascii_letters(text) and not contains_japanese(text)
    )

    if expected_language == "ja" and ja_count == 0 and en_count > 0:
        raise ValueError(
            "Expected slide language 'ja' but slide text looks non-Japanese."
        )
    if expected_language == "en" and ja_count > 0:
        raise ValueError(
            "Expected slide language 'en' but slide text includes Japanese."
        )


def resolve_base_dir(raw_base: str | None) -> Path:
    if raw_base:
        return Path(raw_base).expanduser().resolve()
    return Path.cwd().resolve()


def resolve_spec_input(raw_spec: str, base_dir: Path) -> tuple[dict[str, Any], Path]:
    source = raw_spec.strip()
    if not source:
        raise ValueError("--spec is empty.")

    if source == "-":
        raise ValueError(
            "stdin input is disabled for --spec. Write JSON to a file and pass its path."
        )

    if source.startswith("@"):
        source = source[1:].strip()
        if not source:
            raise ValueError("Invalid --spec value: @ requires a file path.")

    if looks_like_json_object(source):
        raise ValueError(
            "Inline JSON via --spec is disabled to avoid argv size limits. "
            "Write JSON to a file and pass the file path."
        )

    spec_path = Path(source).expanduser().resolve()
    if not spec_path.exists():
        raise FileNotFoundError(
            f"Spec not found: {spec_path}. Provide an existing JSON file path."
        )
    if spec_path.suffix.lower() in {".md", ".ppt", ".pptx"}:
        raise ValueError(
            f"Unsupported spec file type: {spec_path.suffix}. --spec expects slide JSON, not {spec_path.name}."
        )
    return load_json(spec_path), spec_path


def append_title_slide(
    prs: Presentation,
    layouts: list[dict[str, Any]],
    raw_spec: dict[str, Any],
    fallback_layout: int,
    prefer_master_layouts: bool,
) -> tuple[bool, bool, bool, dict[str, Any]]:
    requested_layout_name = str(raw_spec.get("layout_name", "")).strip()
    layout, selection_mode, selection_detail = pick_layout(
        prs,
        layouts,
        raw_spec.get("layout"),
        requested_layout_name,
        fallback_layout,
        prefer_master_layouts=prefer_master_layouts,
        is_title_slide=True,
        has_bullets=False,
        has_table=False,
        has_visual=False,
        has_subtitle=bool(str(raw_spec.get("subtitle", "")).strip()),
        requested_intent=str(raw_spec.get("intent", "")).strip().lower() or None,
    )

    slide = prs.slides.add_slide(layout)
    set_title(slide, str(raw_spec.get("title", "")).strip())
    set_subtitle(slide, str(raw_spec.get("subtitle", "")).strip())

    title_sources = raw_spec.get("sources", [])
    title_notes = raw_spec.get("speaker_notes", [])
    if not isinstance(title_sources, list):
        title_sources = []
    if not isinstance(title_notes, list):
        title_notes = []
    add_notes(slide, [str(item) for item in title_sources], [str(item) for item in title_notes])

    used_master_auto = selection_mode == "master_auto"
    return used_master_auto, False, False, selection_detail


def append_content_slide(
    prs: Presentation,
    layouts: list[dict[str, Any]],
    raw_spec: dict[str, Any],
    spec_reference_path: Path,
    fallback_layout: int,
    strict_images: bool,
    prefer_master_layouts: bool,
) -> tuple[bool, bool, bool, dict[str, Any]]:
    spec = normalize_slide_spec(raw_spec, 1)
    layout, selection_mode, selection_detail = pick_layout(
        prs,
        layouts,
        spec["layout"],
        spec["layout_name"],
        fallback_layout,
        prefer_master_layouts=prefer_master_layouts,
        is_title_slide=False,
        has_bullets=bool(spec["bullets"]),
        has_table=bool(spec["table"]),
        has_visual=bool(spec["visual"].get("path")),
        has_subtitle=bool(spec["subtitle"]),
        bullet_count=len(spec["bullets"]),
        requested_intent=spec.get("intent") or None,
    )

    slide = prs.slides.add_slide(layout)
    set_title(slide, spec["title"])
    set_subtitle(slide, spec["subtitle"])
    reserved_shape_ids: set[int] = set()
    intent_key = str(selection_detail.get("intent_key", "")).strip().lower()
    operation_order = ["visual", "table"]
    if intent_key in {"table_comparison", "table_focus"}:
        operation_order = ["table", "visual"]

    table_added = False
    visual_added = False
    for operation in operation_order:
        if operation == "visual" and spec["visual"].get("path"):
            visual_added = add_visual(
                prs,
                slide,
                spec["visual"],
                spec_reference_path,
                strict_images,
                reserved_shape_ids=reserved_shape_ids,
                has_bullets=bool(spec["bullets"]),
            )
        if operation == "table" and spec["table"]:
            table_added = add_table(
                prs,
                slide,
                spec["table"],
                reserved_shape_ids=reserved_shape_ids,
            )

    set_bullets(slide, spec["bullets"], reserved_shape_ids=reserved_shape_ids)
    add_notes(slide, spec["sources"], spec["speaker_notes"])

    used_master_auto = selection_mode == "master_auto"
    return used_master_auto, table_added, visual_added, selection_detail


def run(args: argparse.Namespace) -> dict[str, Any]:
    deck_path = Path(args.deck).expanduser().resolve()
    if not deck_path.exists():
        raise FileNotFoundError(f"Deck not found: {deck_path}")

    base_dir = resolve_base_dir(args.plan_base)
    raw_spec, spec_reference_path = resolve_spec_input(args.spec, base_dir)
    if not isinstance(raw_spec, dict):
        raise ValueError("Slide spec must be a JSON object.")
    expected_language = normalize_language(args.language)
    if args.language and expected_language is None:
        raise ValueError("Unsupported --language value. Use 'ja' or 'en'.")
    validate_slide_language(raw_spec, expected_language)

    prs = Presentation(str(deck_path))
    layouts = list_master_layouts(prs)
    prefer_master_layouts = not args.no_prefer_master_layouts

    if args.kind == "title":
        used_master_auto, table_added, visual_added, selection_detail = append_title_slide(
            prs,
            layouts,
            raw_spec,
            args.fallback_layout,
            prefer_master_layouts,
        )
    else:
        used_master_auto, table_added, visual_added, selection_detail = append_content_slide(
            prs,
            layouts,
            raw_spec,
            spec_reference_path,
            args.fallback_layout,
            args.strict_images,
            prefer_master_layouts,
        )

    prs.save(str(deck_path))
    return {
        "deck_path": str(deck_path),
        "kind": args.kind,
        "language": expected_language or "auto",
        "used_master_auto": used_master_auto,
        "table_added": table_added,
        "visual_added": visual_added,
        "layout_intent": selection_detail,
    }


def main() -> int:
    args = parse_args()
    try:
        result = run(args)
    except Exception as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 1

    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
