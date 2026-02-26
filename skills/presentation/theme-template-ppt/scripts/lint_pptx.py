#!/usr/bin/env python3
"""Lint a PPTX for density, structure, and visual-design consistency."""

from __future__ import annotations

import argparse
import json
import re
import statistics
import sys
from collections import Counter
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_INCH = 914400
THEME_DEFAULT_FONT = "__theme_default__"
JAPANESE_CHAR_RE = re.compile(r"[ぁ-ゖァ-ヺー一-龯々〆〤]")
ASCII_LETTER_RE = re.compile(r"[A-Za-z]")
NON_CONTENT_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "DATE", None),
        getattr(PP_PLACEHOLDER, "FOOTER", None),
        getattr(PP_PLACEHOLDER, "HEADER", None),
        getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
    )
    if value is not None
}
VISUAL_PLACEHOLDER_TYPES = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "PICTURE", None),
        getattr(PP_PLACEHOLDER, "CHART", None),
        getattr(PP_PLACEHOLDER, "MEDIA_CLIP", None),
        getattr(PP_PLACEHOLDER, "BITMAP", None),
        getattr(PP_PLACEHOLDER, "ORG_CHART", None),
        getattr(PP_PLACEHOLDER, "SLIDE_IMAGE", None),
    )
    if value is not None
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, help="Input PPTX file")
    parser.add_argument(
        "--max-chars",
        type=int,
        default=220,
        help="Maximum body text characters per slide",
    )
    parser.add_argument(
        "--max-bullets",
        type=int,
        default=5,
        help="Maximum number of body bullet lines per slide",
    )
    parser.add_argument(
        "--max-line-chars",
        type=int,
        default=80,
        help="Maximum characters for a single body line",
    )
    parser.add_argument(
        "--min-visual-ratio",
        type=float,
        default=0.50,
        help="Minimum ratio of content slides that contain visuals",
    )
    parser.add_argument(
        "--max-consecutive-text-only",
        type=int,
        default=1,
        help="Maximum consecutive content slides with no visual",
    )
    parser.add_argument(
        "--min-structured-ratio",
        type=float,
        default=0.90,
        help="Minimum ratio of content slides using bullets/table/visual",
    )
    parser.add_argument(
        "--max-unstructured-content-slides",
        type=int,
        default=0,
        help="Maximum count of content slides with no bullets/table/visual",
    )
    parser.add_argument(
        "--max-font-families",
        type=int,
        default=2,
        help="Maximum distinct font families allowed per slide/deck",
    )
    parser.add_argument(
        "--max-font-sizes",
        type=int,
        default=6,
        help="Maximum distinct font sizes allowed per slide/deck",
    )
    parser.add_argument(
        "--min-title-body-size-ratio",
        type=float,
        default=1.15,
        help="Minimum median title/body font-size ratio for hierarchy",
    )
    parser.add_argument(
        "--max-layout-share",
        type=float,
        default=0.65,
        help="Maximum share for the most-used layout among content slides",
    )
    parser.add_argument(
        "--max-image-aspect-ratio-span",
        type=float,
        default=2.5,
        help="Maximum max/min image aspect ratio span (>=3 images)",
    )
    parser.add_argument(
        "--max-left-guides",
        type=int,
        default=4,
        help="Maximum distinct left-guide positions per slide (0.1-inch grid)",
    )
    parser.add_argument(
        "--min-content-margin-inch",
        type=float,
        default=0.30,
        help="Minimum left/right margin for content shapes in inches",
    )
    parser.add_argument(
        "--language",
        choices=["ja", "en"],
        help="Expected deck language. When set, lint warns on mixed-language slides.",
    )
    parser.add_argument(
        "--require-sources-notes",
        action="store_true",
        help="Require 'Sources:' notes on all non-cover slides with body text",
    )
    parser.add_argument(
        "--sources-prefix",
        default="Sources:",
        help="Prefix used in notes to identify source sections",
    )
    parser.add_argument(
        "--json-output",
        help="Optional path to write a JSON report",
    )
    parser.add_argument(
        "--fail-on-warning",
        action="store_true",
        help="Exit with code 1 when warnings are detected",
    )
    return parser.parse_args()


def shape_has_visual(shape: Any) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    placeholder_type = get_placeholder_type(shape)
    if placeholder_type in VISUAL_PLACEHOLDER_TYPES:
        return True
    try:
        if getattr(shape, "has_chart", False) and shape.has_chart:
            return True
    except Exception:
        pass
    return False


def shape_has_table(shape: Any) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.TABLE


def shape_has_non_empty_text(shape: Any) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.text.strip():
            return True
    return False


def get_placeholder_type(shape: Any) -> Any | None:
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return shape.placeholder_format.type
    except Exception:
        return None


def is_non_content_placeholder(shape: Any) -> bool:
    return get_placeholder_type(shape) in NON_CONTENT_PLACEHOLDERS


def contains_japanese(text: str) -> bool:
    return bool(JAPANESE_CHAR_RE.search(text))


def contains_ascii_letters(text: str) -> bool:
    return bool(ASCII_LETTER_RE.search(text))


def normalize_font_name(raw: Any) -> str:
    value = str(raw or "").strip()
    if not value:
        return THEME_DEFAULT_FONT
    return value.lower()


def normalize_font_size_pt(raw: Any) -> float | None:
    if raw is None:
        return None
    try:
        size_pt = float(raw.pt)
    except Exception:
        return None
    return round(size_pt, 1)


def collect_paragraph_style(paragraph: Any) -> tuple[set[str], list[float]]:
    text = paragraph.text.strip()
    if not text:
        return set(), []

    families: set[str] = set()
    sizes: list[float] = []

    paragraph_font = getattr(paragraph, "font", None)
    paragraph_family = normalize_font_name(getattr(paragraph_font, "name", None))
    if paragraph_family:
        families.add(paragraph_family)
    paragraph_size = normalize_font_size_pt(getattr(paragraph_font, "size", None))
    if paragraph_size is not None:
        sizes.append(paragraph_size)

    found_non_empty_run = False
    for run in paragraph.runs:
        if not run.text.strip():
            continue
        found_non_empty_run = True
        run_family = normalize_font_name(getattr(run.font, "name", None))
        if run_family:
            families.add(run_family)
        run_size = normalize_font_size_pt(getattr(run.font, "size", None))
        if run_size is not None:
            sizes.append(run_size)

    if not families and found_non_empty_run:
        families.add(THEME_DEFAULT_FONT)
    if not families:
        families.add(THEME_DEFAULT_FONT)
    return families, sizes


def detect_language_warning(texts: list[str], expected_language: str | None) -> str | None:
    if expected_language is None or not texts:
        return None

    merged = "\n".join(texts)
    has_ja = contains_japanese(merged)
    has_en = contains_ascii_letters(merged)
    if expected_language == "ja" and not has_ja and has_en:
        return "expected Japanese text, but slide looks non-Japanese"
    if expected_language == "en" and has_ja:
        return "expected English text, but slide includes Japanese"
    return None


def quantize_guide(value_emu: int) -> float:
    inches = float(value_emu) / EMU_PER_INCH
    return round(inches, 1)


def collect_slide_metrics(
    slide: Any,
    index: int,
    *,
    slide_width: int,
    sources_prefix: str,
    expected_language: str | None,
) -> dict[str, Any]:
    title_shape = slide.shapes.title
    has_title = False
    body_chars = 0
    body_lines = 0
    longest_line = 0
    has_visual = False
    has_table = False
    all_texts: list[str] = []

    font_families: set[str] = set()
    font_sizes: set[float] = set()
    title_font_sizes: list[float] = []
    body_font_sizes: list[float] = []
    image_aspect_ratios: list[float] = []
    left_guides: set[float] = set()
    content_shape_count = 0
    min_content_margin_inch: float | None = None

    layout_name = ""
    try:
        layout_name = str(getattr(slide.slide_layout, "name", "") or "").strip()
    except Exception:
        layout_name = ""

    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        has_title = bool(title_shape.text.strip())

    for shape in slide.shapes:
        if is_non_content_placeholder(shape):
            continue
        is_title_shape = title_shape is not None and shape == title_shape

        if shape_has_visual(shape):
            has_visual = True
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and getattr(shape, "height", 0):
                image_aspect_ratios.append(round(float(shape.width) / float(shape.height), 4))
        if shape_has_table(shape):
            has_table = True

        is_content_shape = (
            (not is_title_shape and getattr(shape, "has_text_frame", False))
            or shape_has_visual(shape)
            or shape_has_table(shape)
        )
        if is_content_shape:
            content_shape_count += 1
            left_guides.add(quantize_guide(int(getattr(shape, "left", 0))))
            left_margin = float(getattr(shape, "left", 0)) / EMU_PER_INCH
            right_edge = float(getattr(shape, "left", 0) + getattr(shape, "width", 0))
            right_margin = float(slide_width - right_edge) / EMU_PER_INCH
            local_min = min(left_margin, right_margin)
            if min_content_margin_inch is None:
                min_content_margin_inch = local_min
            else:
                min_content_margin_inch = min(min_content_margin_inch, local_min)

        if not getattr(shape, "has_text_frame", False):
            continue

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            all_texts.append(text)
            families, sizes = collect_paragraph_style(paragraph)
            font_families.update(families)
            font_sizes.update(sizes)

            if is_title_shape:
                title_font_sizes.extend(sizes)
                continue

            body_lines += 1
            body_chars += len(text)
            longest_line = max(longest_line, len(text))
            body_font_sizes.extend(sizes)

    has_bullets = body_lines > 0
    has_structured = has_bullets or has_table or has_visual
    is_cover_like = has_title and body_lines <= 1 and not has_table and not has_visual

    notes_text = ""
    if getattr(slide, "has_notes_slide", False):
        try:
            notes_text = str(slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
    has_sources_note = sources_prefix.lower() in notes_text.lower()

    language_warning = detect_language_warning(all_texts, expected_language)

    return {
        "slide": index,
        "layout_name": layout_name,
        "has_title": has_title,
        "body_chars": body_chars,
        "body_lines": body_lines,
        "longest_line": longest_line,
        "has_bullets": has_bullets,
        "has_table": has_table,
        "has_visual": has_visual,
        "has_structured": has_structured,
        "is_cover_like": is_cover_like,
        "font_families": sorted(font_families),
        "font_sizes": sorted(font_sizes),
        "title_font_sizes": [round(value, 1) for value in title_font_sizes],
        "body_font_sizes": [round(value, 1) for value in body_font_sizes],
        "left_guides": sorted(left_guides),
        "left_guide_count": len(left_guides),
        "content_shape_count": content_shape_count,
        "min_content_margin_inch": None
        if min_content_margin_inch is None
        else round(min_content_margin_inch, 3),
        "image_aspect_ratios": image_aspect_ratios,
        "has_sources_note": has_sources_note,
        "notes_present": bool(notes_text),
        "language_warning": language_warning,
        "warnings": [],
    }


def lint(args: argparse.Namespace) -> dict[str, Any]:
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"PPTX not found: {input_path}")

    prs = Presentation(str(input_path))
    slides = list(prs.slides)
    if not slides:
        raise ValueError("Deck has zero slides.")

    results = [
        collect_slide_metrics(
            slide,
            idx,
            slide_width=int(prs.slide_width),
            sources_prefix=args.sources_prefix,
            expected_language=args.language,
        )
        for idx, slide in enumerate(slides, start=1)
    ]

    for item in results:
        if not item["has_title"]:
            item["warnings"].append("missing title text")
        if item["body_chars"] > args.max_chars:
            item["warnings"].append(
                f"body chars {item['body_chars']} > max {args.max_chars}"
            )
        if item["body_lines"] > args.max_bullets:
            item["warnings"].append(
                f"body lines {item['body_lines']} > max {args.max_bullets}"
            )
        if item["longest_line"] > args.max_line_chars:
            item["warnings"].append(
                f"line chars {item['longest_line']} > max {args.max_line_chars}"
            )
        if not item["has_structured"] and not item["is_cover_like"]:
            item["warnings"].append("unstructured content; add bullets, table, or visual")
        if item["left_guide_count"] > args.max_left_guides and not item["is_cover_like"]:
            item["warnings"].append(
                "too many horizontal guide positions; align content to fewer columns"
            )
        min_margin = item["min_content_margin_inch"]
        if min_margin is not None and min_margin < args.min_content_margin_inch and not item[
            "is_cover_like"
        ]:
            item["warnings"].append(
                "content too close to edge; increase left/right margin consistency"
            )
        if len(item["font_families"]) > args.max_font_families:
            item["warnings"].append(
                f"font families {len(item['font_families'])} > max {args.max_font_families}"
            )
        if len(item["font_sizes"]) > args.max_font_sizes:
            item["warnings"].append(
                f"font sizes {len(item['font_sizes'])} > max {args.max_font_sizes}"
            )
        if args.require_sources_notes and not item["is_cover_like"] and item["body_chars"] > 0:
            if not item["has_sources_note"]:
                item["warnings"].append(
                    f"missing notes section '{args.sources_prefix}'"
                )
        if item["language_warning"]:
            item["warnings"].append(item["language_warning"])

    content_slides = [item for item in results if not item["is_cover_like"]]
    content_slide_count = len(content_slides)

    streak = 0
    for item in results:
        if item["is_cover_like"]:
            streak = 0
            continue
        if item["has_visual"] or item["has_table"]:
            streak = 0
            continue
        streak += 1
        if streak > args.max_consecutive_text_only:
            item["warnings"].append(
                "too many consecutive text-only slides; add chart/diagram/table/image"
            )

    total_slides = len(results)
    visual_slides = sum(1 for item in results if item["has_visual"])
    table_slides = sum(1 for item in results if item["has_table"])
    structured_slides = sum(1 for item in results if item["has_structured"])
    visual_content_slides = sum(1 for item in content_slides if item["has_visual"])
    unstructured_content_slides = sum(1 for item in content_slides if not item["has_structured"])
    visual_ratio = (
        1.0 if content_slide_count == 0 else visual_content_slides / content_slide_count
    )
    structured_ratio = (
        1.0
        if content_slide_count == 0
        else (content_slide_count - unstructured_content_slides) / content_slide_count
    )

    layout_counts: Counter[str] = Counter(
        (item["layout_name"] or "<unknown>") for item in content_slides
    )
    top_layout = ""
    top_layout_count = 0
    max_layout_share = 0.0
    if layout_counts and content_slide_count > 0:
        top_layout, top_layout_count = layout_counts.most_common(1)[0]
        max_layout_share = top_layout_count / content_slide_count

    deck_font_families = sorted(
        {family for item in results for family in item["font_families"]}
    )
    deck_font_sizes = sorted(
        {size for item in results for size in item["font_sizes"]}
    )
    title_sizes = [size for item in results for size in item["title_font_sizes"]]
    body_sizes = [size for item in results for size in item["body_font_sizes"]]
    title_body_size_ratio: float | None = None
    if title_sizes and body_sizes:
        median_body = statistics.median(body_sizes)
        if median_body > 0:
            title_body_size_ratio = statistics.median(title_sizes) / median_body

    image_ratios = [ratio for item in results for ratio in item["image_aspect_ratios"]]
    image_aspect_ratio_span: float | None = None
    if len(image_ratios) >= 3:
        min_ratio = min(image_ratios)
        max_ratio = max(image_ratios)
        if min_ratio > 0:
            image_aspect_ratio_span = max_ratio / min_ratio

    source_note_coverage = 1.0
    if content_slide_count > 0:
        source_note_coverage = (
            sum(1 for item in content_slides if item["has_sources_note"])
            / content_slide_count
        )

    deck_warnings: list[str] = []
    if visual_ratio < args.min_visual_ratio:
        deck_warnings.append(
            f"visual ratio {visual_ratio:.2f} < min {args.min_visual_ratio:.2f}"
        )
    if structured_ratio < args.min_structured_ratio:
        deck_warnings.append(
            f"structured ratio {structured_ratio:.2f} < min {args.min_structured_ratio:.2f}"
        )
    if unstructured_content_slides > args.max_unstructured_content_slides:
        deck_warnings.append(
            "unstructured content slides "
            f"{unstructured_content_slides} > max {args.max_unstructured_content_slides}"
        )
    if max_layout_share > args.max_layout_share:
        deck_warnings.append(
            f"layout usage biased: '{top_layout}' share {max_layout_share:.2f} > max {args.max_layout_share:.2f}"
        )
    if len(deck_font_families) > args.max_font_families:
        deck_warnings.append(
            f"deck font families {len(deck_font_families)} > max {args.max_font_families}"
        )
    if len(deck_font_sizes) > args.max_font_sizes:
        deck_warnings.append(
            f"deck font sizes {len(deck_font_sizes)} > max {args.max_font_sizes}"
        )
    if (
        title_body_size_ratio is not None
        and title_body_size_ratio < args.min_title_body_size_ratio
    ):
        deck_warnings.append(
            "title/body size hierarchy too weak: "
            f"{title_body_size_ratio:.2f} < min {args.min_title_body_size_ratio:.2f}"
        )
    if (
        image_aspect_ratio_span is not None
        and image_aspect_ratio_span > args.max_image_aspect_ratio_span
    ):
        deck_warnings.append(
            "image aspect ratios vary too much: "
            f"{image_aspect_ratio_span:.2f} > max {args.max_image_aspect_ratio_span:.2f}"
        )
    if args.require_sources_notes and source_note_coverage < 1.0:
        deck_warnings.append(
            f"source note coverage {source_note_coverage:.2f} < 1.00 on content slides"
        )

    warning_count = len(deck_warnings) + sum(len(item["warnings"]) for item in results)
    report = {
        "input": str(input_path),
        "summary": {
            "slides": total_slides,
            "content_slides": content_slide_count,
            "visual_slides": visual_slides,
            "visual_content_slides": visual_content_slides,
            "table_slides": table_slides,
            "structured_slides": structured_slides,
            "visual_ratio": round(visual_ratio, 4),
            "structured_ratio": round(structured_ratio, 4),
            "unstructured_content_slides": unstructured_content_slides,
            "distinct_layouts": len(layout_counts),
            "top_layout": top_layout,
            "top_layout_count": top_layout_count,
            "max_layout_share": round(max_layout_share, 4),
            "deck_font_families": len(deck_font_families),
            "deck_font_sizes": len(deck_font_sizes),
            "title_body_size_ratio": None
            if title_body_size_ratio is None
            else round(title_body_size_ratio, 4),
            "image_count": len(image_ratios),
            "image_aspect_ratio_span": None
            if image_aspect_ratio_span is None
            else round(image_aspect_ratio_span, 4),
            "source_note_coverage": round(source_note_coverage, 4),
            "warning_count": warning_count,
        },
        "deck_warnings": deck_warnings,
        "slides": results,
    }
    return report


def print_report(report: dict[str, Any]) -> None:
    summary = report["summary"]
    print(f"Deck: {report['input']}")
    print(
        "Summary: slides={slides}, content_slides={content_slides}, visual_ratio={visual_ratio:.2f}, structured_ratio={structured_ratio:.2f}, top_layout='{top_layout}' ({max_layout_share:.2f}), fonts={deck_font_families}, sizes={deck_font_sizes}, warnings={warning_count}".format(
            **summary
        )
    )

    for warning in report["deck_warnings"]:
        print(f"Deck warning: {warning}")

    for item in report["slides"]:
        visual = "yes" if item["has_visual"] else "no"
        table = "yes" if item["has_table"] else "no"
        bullets = "yes" if item["has_bullets"] else "no"
        status = "WARN" if item["warnings"] else "OK"
        print(
            "Slide {slide:02d}: status={status}, layout={layout}, body_chars={chars}, body_lines={lines}, guides={guides}, margin={margin}, bullets={bullets}, table={table}, visual={visual}".format(
                slide=item["slide"],
                status=status,
                layout=item["layout_name"] or "<unknown>",
                chars=item["body_chars"],
                lines=item["body_lines"],
                guides=item["left_guide_count"],
                margin=item["min_content_margin_inch"],
                bullets=bullets,
                table=table,
                visual=visual,
            )
        )
        for warning in item["warnings"]:
            print(f"  - {warning}")


def write_json_report(path: str, report: dict[str, Any]) -> None:
    output_path = Path(path).expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps(report, indent=2), encoding="utf-8")
    print(f"[OK] Wrote JSON report: {output_path}")


def main() -> int:
    args = parse_args()
    try:
        report = lint(args)
    except Exception as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 1

    print_report(report)

    if args.json_output:
        write_json_report(args.json_output, report)

    has_warnings = bool(report["deck_warnings"]) or any(
        item["warnings"] for item in report["slides"]
    )
    if args.fail_on_warning and has_warnings:
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
