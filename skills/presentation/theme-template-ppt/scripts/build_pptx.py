#!/usr/bin/env python3
"""Build a PowerPoint file from a template and a deck-plan JSON."""

from __future__ import annotations

import argparse
import json
import os
import posixpath
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches, Pt


BODY_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "BODY", None),
        getattr(PP_PLACEHOLDER, "OBJECT", None),
        getattr(PP_PLACEHOLDER, "TEXT", None),
        getattr(PP_PLACEHOLDER, "VERTICAL_BODY", None),
        getattr(PP_PLACEHOLDER, "VERTICAL_OBJECT", None),
    )
    if value is not None
}

TITLE_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "TITLE", None),
        getattr(PP_PLACEHOLDER, "CENTER_TITLE", None),
    )
    if value is not None
}

SUBTITLE_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "SUBTITLE", None),
    )
    if value is not None
}

VISUAL_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "PICTURE", None),
        getattr(PP_PLACEHOLDER, "CHART", None),
        getattr(PP_PLACEHOLDER, "MEDIA_CLIP", None),
        getattr(PP_PLACEHOLDER, "BITMAP", None),
        getattr(PP_PLACEHOLDER, "ORG_CHART", None),
        getattr(PP_PLACEHOLDER, "SLIDE_IMAGE", None),
    )
    if value is not None
}

TABLE_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "TABLE", None),
    )
    if value is not None
}

NON_CONTENT_PLACEHOLDERS = {
    value
    for value in (
        getattr(PP_PLACEHOLDER, "DATE", None),
        getattr(PP_PLACEHOLDER, "FOOTER", None),
        getattr(PP_PLACEHOLDER, "HEADER", None),
        getattr(PP_PLACEHOLDER, "SLIDE_NUMBER", None),
    )
    if value is not None
}

VISUAL_LAYOUT_KEYWORDS = (
    "picture",
    "photo",
    "image",
    "media",
    "chart",
    "diagram",
    "caption",
    "two content",
)

TABLE_LAYOUT_KEYWORDS = (
    "table",
    "matrix",
    "comparison",
)

COVER_LAYOUT_KEYWORDS = (
    "title slide",
    "cover",
)

SECTION_LAYOUT_KEYWORDS = (
    "section",
)

AUTO_SLIDE_ORDER_DEFAULT = ("agenda", "title", "content", "summary")
AUTO_SLIDE_ORDER_ALIASES = {
    "agenda": "agenda",
    "toc": "agenda",
    "table_of_contents": "agenda",
    "目次": "agenda",
    "title": "title",
    "title_slide": "title",
    "cover": "title",
    "content": "content",
    "slides": "content",
    "body": "content",
    "summary": "summary",
    "wrapup": "summary",
    "まとめ": "summary",
}

JAPANESE_CHAR_RE = re.compile(r"[ぁ-ゖァ-ヺー一-龯々〆〤]")
ASCII_LETTER_RE = re.compile(r"[A-Za-z]")

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
SLIDE_LAYOUT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
SLIDE_LAYOUT_CONTENT_TYPE = (
    "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"
)
FIXED_WORK_ROOT = Path("~/.foundry_local_playground/outputs/pptx").expanduser()
WORKDIR_KEY_ENV = "THEME_TEMPLATE_PPT_WORK_KEY"
WORKDIR_KEY_ENV_FALLBACKS = (
    "PPTX_WORKDIR_KEY",
    "CODEX_THREAD_ID",
)


def normalize_language(raw: Any) -> str | None:
    if raw is None:
        return None
    value = str(raw).strip().lower()
    if value in {"ja", "jp", "japanese"}:
        return "ja"
    if value in {"en", "english"}:
        return "en"
    return None


def contains_japanese(text: str) -> bool:
    return bool(JAPANESE_CHAR_RE.search(text))


def contains_ascii_letters(text: str) -> bool:
    return bool(ASCII_LETTER_RE.search(text))


def infer_language_from_plan(plan: dict[str, Any], slides_raw: list[Any]) -> str:
    explicit = normalize_language(plan.get("language"))
    if explicit:
        return explicit

    texts: list[str] = []
    title_slide = plan.get("title_slide")
    if isinstance(title_slide, dict):
        for key in ("title", "subtitle"):
            value = str(title_slide.get(key, "")).strip()
            if value:
                texts.append(value)

    for raw_slide in slides_raw:
        if not isinstance(raw_slide, dict):
            continue
        for key in ("title", "subtitle"):
            value = str(raw_slide.get(key, "")).strip()
            if value:
                texts.append(value)
        bullets_raw = raw_slide.get("bullets", [])
        if isinstance(bullets_raw, list):
            for item in bullets_raw:
                value = str(item).strip()
                if value:
                    texts.append(value)

    ja_count = sum(1 for item in texts if contains_japanese(item))
    en_count = sum(
        1
        for item in texts
        if contains_ascii_letters(item) and not contains_japanese(item)
    )
    if ja_count > 0:
        return "ja"
    if en_count > 0:
        return "en"
    return "en"


def extract_content_titles(slides_raw: list[Any]) -> list[str]:
    titles: list[str] = []
    for index, raw_slide in enumerate(slides_raw, start=1):
        if not isinstance(raw_slide, dict):
            titles.append(f"Slide {index}")
            continue
        title = str(raw_slide.get("title", "")).strip()
        titles.append(title or f"Slide {index}")
    return titles


def build_agenda_slide(slides_raw: list[Any], language: str) -> dict[str, Any]:
    titles = extract_content_titles(slides_raw)
    max_items = 8
    bullets = [f"{idx}. {title}" for idx, title in enumerate(titles[:max_items], start=1)]
    if len(titles) > max_items:
        remainder = len(titles) - max_items
        if language == "ja":
            bullets.append(f"他 {remainder} 項目")
        else:
            bullets.append(f"... and {remainder} more")

    if language == "ja":
        return {
            "title": "目次",
            "subtitle": "この資料で扱う項目です",
            "bullets": bullets,
            "speaker_notes": ["自動追加: 先頭目次スライド"],
        }
    return {
        "title": "Agenda",
        "subtitle": "Topics covered in this deck",
        "bullets": bullets,
        "speaker_notes": ["Auto-added: first agenda slide"],
    }


def build_summary_slide(slides_raw: list[Any], language: str) -> dict[str, Any]:
    titles = extract_content_titles(slides_raw)
    topic_count = len(titles)
    top_topics = " / ".join(titles[:3]) if titles else ""

    if language == "ja":
        bullets = [f"全 {topic_count} 項目の要点を整理しました"]
        if top_topics:
            bullets.append(f"主要トピック: {top_topics}")
        bullets.append("次のアクションと意思決定項目を確認してください")
        return {
            "title": "まとめ",
            "subtitle": "本資料の結論",
            "bullets": bullets,
            "speaker_notes": ["自動追加: 末尾まとめスライド"],
        }

    bullets = [f"This deck summarized {topic_count} key topics"]
    if top_topics:
        bullets.append(f"Main topics: {top_topics}")
    bullets.append("Confirm the next actions and decisions")
    return {
        "title": "Summary",
        "subtitle": "Key takeaways",
        "bullets": bullets,
        "speaker_notes": ["Auto-added: final summary slide"],
    }


def normalize_auto_slide_order(raw: Any) -> list[str]:
    if raw is None:
        return list(AUTO_SLIDE_ORDER_DEFAULT)

    if isinstance(raw, str):
        tokens_raw = [item.strip() for item in raw.split(",") if item.strip()]
    elif isinstance(raw, list):
        tokens_raw = [str(item).strip() for item in raw if str(item).strip()]
    else:
        raise ValueError("Field 'auto_slide_order' must be a list or comma-separated string.")

    order: list[str] = []
    seen: set[str] = set()
    for token in tokens_raw:
        mapped = AUTO_SLIDE_ORDER_ALIASES.get(token.lower()) or AUTO_SLIDE_ORDER_ALIASES.get(token)
        if mapped is None:
            allowed = ", ".join(sorted(set(AUTO_SLIDE_ORDER_ALIASES.keys())))
            raise ValueError(
                f"Unsupported auto_slide_order token: {token!r}. Allowed aliases: {allowed}"
            )
        if mapped in seen:
            continue
        seen.add(mapped)
        order.append(mapped)

    for token in AUTO_SLIDE_ORDER_DEFAULT:
        if token not in seen:
            order.append(token)
    return order


def slide_requests_visual(raw_slide: dict[str, Any]) -> bool:
    visual = raw_slide.get("visual")
    if not isinstance(visual, dict):
        return False
    return bool(str(visual.get("path", "")).strip())


def slide_requests_table(raw_slide: dict[str, Any]) -> bool:
    table = raw_slide.get("table")
    if not isinstance(table, dict):
        return False
    headers = table.get("headers")
    rows = table.get("rows")
    return bool(isinstance(headers, list) and headers and isinstance(rows, list) and rows)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--slide-title", help="Optional deck title used for output file naming")
    parser.add_argument(
        "--template",
        help="Path to PPTX template (default: <skill-dir>/assets/template.pptx)",
    )
    parser.add_argument(
        "--plan",
        help="Deck plan JSON file path",
    )
    parser.add_argument(
        "--work-root",
        default=None,
        help="Deprecated (ignored). Work dir is fixed to ~/.foundry_local_playground/outputs/pptx/<work_key>/",
    )
    parser.add_argument(
        "--thread-key",
        help=(
            "Optional thread identifier used to reuse one work directory in the same thread "
            f"(default: {WORKDIR_KEY_ENV}, fallback: {', '.join(WORKDIR_KEY_ENV_FALLBACKS)})."
        ),
    )
    parser.add_argument(
        "--output",
        help="Optional output file name (.pptx). Saved inside workdir.",
    )
    parser.add_argument(
        "--fallback-layout",
        type=int,
        default=1,
        help="Deprecated (ignored). Layout is always auto-detected from slide masters.",
    )
    parser.add_argument(
        "--strict-images",
        action="store_true",
        help="Fail when a visual image path does not exist",
    )
    parser.add_argument(
        "--no-prefer-master-layouts",
        action="store_true",
        help="Deprecated (ignored). Master layout auto-detection is always enabled.",
    )
    parser.add_argument(
        "--list-layouts",
        action="store_true",
        help=(
            "Inspect template slide masters and print a JSON layout catalog. "
            "When used, --plan is not required."
        ),
    )
    parser.add_argument(
        "--layout-report",
        help="Optional JSON output file for --list-layouts mode",
    )

    args = parser.parse_args()
    if args.list_layouts:
        return args
    if not args.plan:
        parser.error("--plan is required unless --list-layouts is used.")
    return args


def normalize_slide_title(raw: str) -> str:
    cleaned = raw.strip()
    if not cleaned:
        raise ValueError("slide-title is empty.")
    cleaned = re.sub(r"[\\/]+", "-", cleaned)
    cleaned = re.sub(r"\s+", "-", cleaned)
    cleaned = re.sub(r"-{2,}", "-", cleaned).strip("-")
    if not cleaned:
        raise ValueError("slide-title became empty after normalization.")
    return cleaned


def generate_workdir_key() -> str:
    return f"wk-{uuid.uuid4().hex}"


def resolve_thread_key(args: argparse.Namespace) -> tuple[str, str, bool]:
    from_arg = str(args.thread_key or "").strip()
    if from_arg:
        return from_arg, "arg:--thread-key", False

    from_env = str(os.getenv(WORKDIR_KEY_ENV) or "").strip()
    if from_env:
        return from_env, f"env:{WORKDIR_KEY_ENV}", False

    for fallback_env in WORKDIR_KEY_ENV_FALLBACKS:
        fallback_value = str(os.getenv(fallback_env) or "").strip()
        if fallback_value:
            return fallback_value, f"env:{fallback_env}", False

    generated = generate_workdir_key()
    os.environ[WORKDIR_KEY_ENV] = generated
    return generated, "generated", True


def normalize_thread_dir_name(raw: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "-", raw).strip("-._")
    return value or "default-thread"


def resolve_skill_dir() -> Path:
    current = Path(__file__).resolve().parent
    for candidate in [current, *current.parents]:
        if (candidate / "SKILL.md").exists():
            return candidate
    raise RuntimeError(f"Unable to resolve skill directory from script path: {__file__}")


def resolve_template_path(raw_template: str | None) -> Path:
    if raw_template:
        return Path(raw_template).expanduser().resolve()
    return (resolve_skill_dir() / "assets" / "template.pptx").resolve()


def prepare_work_paths(
    args: argparse.Namespace, template_path: Path
) -> tuple[Path, Path, Path, str, str, bool, bool]:
    root = FIXED_WORK_ROOT.resolve()
    root.mkdir(parents=True, exist_ok=True)

    thread_key_raw, thread_key_source, generated_workdir_key = resolve_thread_key(args)
    thread_key = normalize_thread_dir_name(thread_key_raw)
    os.environ[WORKDIR_KEY_ENV] = thread_key
    work_dir_name = normalize_thread_dir_name(thread_key)
    work_dir = (root / work_dir_name).resolve()
    reused_work_dir = work_dir.exists()
    work_dir.mkdir(parents=True, exist_ok=True)

    staged_template = work_dir / template_path.name
    if template_path.resolve() != staged_template.resolve():
        shutil.copy2(template_path, staged_template)

    if args.output:
        output_name = Path(args.output).name
        if not output_name.lower().endswith(".pptx"):
            output_name = f"{output_name}.pptx"
    else:
        if args.slide_title:
            output_name = f"{normalize_slide_title(args.slide_title)}.pptx"
        else:
            output_name = "deck.pptx"
    output_path = work_dir / output_name

    return (
        work_dir,
        staged_template,
        output_path,
        thread_key,
        thread_key_source,
        generated_workdir_key,
        reused_work_dir,
    )


def load_json(path: Path) -> dict[str, Any]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError(
            f"Plan file is not UTF-8 text: {path}. --plan expects a UTF-8 JSON file."
        ) from exc
    return load_json_text(text, str(path))


def load_json_text(raw: str, source_label: str) -> dict[str, Any]:
    try:
        return json.loads(raw)
    except json.JSONDecodeError as exc:
        raise ValueError(f"Invalid JSON: {source_label}: {exc}") from exc


def looks_like_json_object(raw: str) -> bool:
    return raw.lstrip().startswith("{")


def resolve_plan_input(raw_plan: str, work_dir: Path) -> tuple[dict[str, Any], Path, Path]:
    source = raw_plan.strip()
    if not source:
        raise ValueError("--plan is empty.")

    if source == "-":
        raise ValueError(
            "stdin input is disabled for --plan. Write JSON to a file and pass its path."
        )

    if source.startswith("@"):
        source = source[1:].strip()
        if not source:
            raise ValueError("Invalid --plan value: @ requires a file path.")

    if looks_like_json_object(source):
        raise ValueError(
            "Inline JSON via --plan is disabled to avoid argv size limits. "
            "Write JSON to a file and pass the file path."
        )

    plan_path = Path(source).expanduser().resolve()
    if not plan_path.exists():
        raise FileNotFoundError(
            f"Plan not found: {plan_path}. Provide an existing JSON file path via --plan."
        )
    if plan_path.suffix.lower() in {".md", ".ppt", ".pptx"}:
        raise ValueError(
            f"Unsupported plan file type: {plan_path.suffix}. --plan expects deck plan JSON, not {plan_path.name}."
        )

    plan = load_json(plan_path)
    staged_plan = (work_dir / plan_path.name).resolve()
    if plan_path.resolve() != staged_plan:
        shutil.copy2(plan_path, staged_plan)
    return plan, plan_path, staged_plan


def to_emu(value: Any, total: int, default_emu: int) -> int:
    if value is None:
        return default_emu
    if isinstance(value, (int, float)):
        return int(Inches(float(value)))
    if isinstance(value, str):
        raw = value.strip()
        if raw.endswith("%"):
            pct = float(raw[:-1])
            return int(total * pct / 100.0)
        return int(Inches(float(raw)))
    raise ValueError(f"Unsupported dimension value: {value!r}")


def resolve_path(raw: str, plan_path: Path) -> Path:
    path = Path(raw)
    if path.is_absolute():
        return path
    return (plan_path.parent / path).resolve()


def _qn(ns: str, local: str) -> str:
    return f"{{{ns}}}{local}"


def _normalize_rel_target(base_part: str, target: str) -> str:
    base_dir = posixpath.dirname(base_part)
    normalized = posixpath.normpath(posixpath.join(base_dir, target))
    return normalized.lstrip("/")


def _layout_rels_part(layout_part: str) -> str:
    layout_name = posixpath.basename(layout_part)
    return f"ppt/slideLayouts/_rels/{layout_name}.rels"


def _extract_layout_number(layout_part: str) -> int:
    matched = re.search(r"slideLayout(\d+)\.xml$", layout_part)
    if not matched:
        return 0
    return int(matched.group(1))


def _layout_features(layout_root: ET.Element) -> tuple[str, bool, bool, bool]:
    c_sld = layout_root.find(_qn(PML_NS, "cSld"))
    layout_name = ""
    if c_sld is not None:
        layout_name = str(c_sld.attrib.get("name", "")).strip()
    name_lower = layout_name.lower()

    placeholder_types: list[str] = []
    for ph in layout_root.findall(f".//{_qn(PML_NS, 'ph')}"):
        placeholder_types.append(str(ph.attrib.get("type", "")).strip())

    has_title = any(item in {"title", "ctrTitle"} for item in placeholder_types)
    has_subtitle = any(item == "subTitle" for item in placeholder_types)
    has_body = any(
        item in {"", "body", "obj", "tx", "pic", "chart", "tbl", "media", "clipArt"}
        for item in placeholder_types
    )
    return name_lower, has_title, has_subtitle, has_body


def _choose_source_layout(entries: list[dict[str, Any]], profile: str) -> dict[str, Any] | None:
    if not entries:
        return None

    def find_first(predicate: Any) -> dict[str, Any] | None:
        for item in entries:
            if predicate(item):
                return item
        return None

    if profile == "title":
        return (
            find_first(lambda e: e["has_title"] and e["has_subtitle"])
            or find_first(lambda e: e["has_title"])
            or find_first(lambda e: e["has_body"])
            or entries[0]
        )

    return (
        find_first(lambda e: e["has_title"] and e["has_body"])
        or find_first(lambda e: e["has_body"])
        or find_first(lambda e: e["has_title"])
        or entries[0]
    )


def _profile_match(profile: str, entry: dict[str, Any]) -> bool:
    if profile == "title":
        return bool(entry["has_title"] and entry["has_subtitle"])
    if profile == "content":
        return bool(entry["has_title"] and entry["has_body"])
    if profile == "visual":
        return bool(
            entry["has_body"]
            and any(token in entry["name_lower"] for token in VISUAL_LAYOUT_KEYWORDS)
        )
    if profile == "table":
        return bool(
            entry["has_body"]
            and any(token in entry["name_lower"] for token in TABLE_LAYOUT_KEYWORDS)
        )
    return False


def _next_unique_layout_name(base_name: str, used_names: set[str]) -> str:
    candidate = base_name.strip() or "Auto Layout"
    if candidate.lower() not in used_names:
        used_names.add(candidate.lower())
        return candidate
    index = 2
    while True:
        named = f"{candidate} {index}"
        if named.lower() not in used_names:
            used_names.add(named.lower())
            return named
        index += 1


def ensure_required_master_layouts(
    template_path: Path,
    *,
    need_title_layout: bool,
    need_content_layout: bool,
    need_visual_layout: bool,
    need_table_layout: bool,
) -> dict[str, Any]:
    required_profiles: list[tuple[str, str]] = []
    if need_title_layout:
        required_profiles.append(("title", "Auto Title Slide"))
    if need_content_layout:
        required_profiles.append(("content", "Auto Title and Content"))
    if need_visual_layout:
        required_profiles.append(("visual", "Auto Visual Content"))
    if need_table_layout:
        required_profiles.append(("table", "Auto Table Content"))

    if not required_profiles:
        return {"created_count": 0, "created_layout_names": []}

    master_part = "ppt/slideMasters/slideMaster1.xml"
    master_rels_part = "ppt/slideMasters/_rels/slideMaster1.xml.rels"
    content_types_part = "[Content_Types].xml"

    with zipfile.ZipFile(template_path, "r") as zin:
        names = set(zin.namelist())
        if master_part not in names or master_rels_part not in names or content_types_part not in names:
            raise ValueError(
                "Template is missing required OpenXML parts for slide-master layout updates."
            )

        master_root = ET.fromstring(zin.read(master_part))
        master_rels_root = ET.fromstring(zin.read(master_rels_part))
        content_types_root = ET.fromstring(zin.read(content_types_part))

        rel_part_by_id: dict[str, str] = {}
        used_rel_ids: set[str] = set()
        for rel in master_rels_root.findall(_qn(PKG_REL_NS, "Relationship")):
            rel_id = str(rel.attrib.get("Id", "")).strip()
            rel_type = str(rel.attrib.get("Type", "")).strip()
            target = str(rel.attrib.get("Target", "")).strip()
            if rel_id:
                used_rel_ids.add(rel_id)
            if rel_type != SLIDE_LAYOUT_REL_TYPE or not rel_id or not target:
                continue
            rel_part_by_id[rel_id] = _normalize_rel_target(master_part, target)

        layout_id_list = master_root.find(_qn(PML_NS, "sldLayoutIdLst"))
        if layout_id_list is None:
            raise ValueError("slideMaster1.xml does not contain p:sldLayoutIdLst.")

        entries: list[dict[str, Any]] = []
        used_layout_ids: set[int] = set()
        used_layout_name_lowers: set[str] = set()
        max_layout_number = 0
        for order_index, item in enumerate(layout_id_list.findall(_qn(PML_NS, "sldLayoutId"))):
            raw_layout_id = str(item.attrib.get("id", "")).strip()
            if raw_layout_id.isdigit():
                used_layout_ids.add(int(raw_layout_id))
            rel_id = str(item.attrib.get(_qn(REL_NS, "id"), "")).strip()
            if not rel_id:
                continue
            part = rel_part_by_id.get(rel_id)
            if not part or part not in names:
                continue
            max_layout_number = max(max_layout_number, _extract_layout_number(part))
            layout_root = ET.fromstring(zin.read(part))
            name_lower, has_title, has_subtitle, has_body = _layout_features(layout_root)
            used_layout_name_lowers.add(name_lower)
            entries.append(
                {
                    "part": part,
                    "rel_id": rel_id,
                    "name_lower": name_lower,
                    "has_title": has_title,
                    "has_subtitle": has_subtitle,
                    "has_body": has_body,
                    "order_index": order_index,
                }
            )

        if not entries:
            raise ValueError("Template does not contain any slide-master layouts.")

        existing_override_parts = {
            str(override.attrib.get("PartName", "")).lstrip("/")
            for override in content_types_root.findall(_qn(CT_NS, "Override"))
        }

        created_layout_names: list[str] = []
        new_parts: dict[str, bytes] = {}

        next_layout_id = (max(used_layout_ids) + 1) if used_layout_ids else 100000

        def next_rel_id() -> str:
            counter = 1
            while True:
                candidate = f"rId{counter}"
                if candidate not in used_rel_ids:
                    used_rel_ids.add(candidate)
                    return candidate
                counter += 1

        for profile, base_name in required_profiles:
            if any(_profile_match(profile, entry) for entry in entries):
                continue

            source = _choose_source_layout(entries, profile)
            if source is None:
                continue

            max_layout_number += 1
            new_layout_part = f"ppt/slideLayouts/slideLayout{max_layout_number}.xml"
            new_layout_rels_part = _layout_rels_part(new_layout_part)

            source_layout_xml = zin.read(source["part"])
            source_layout_rels_xml = zin.read(_layout_rels_part(source["part"]))
            new_layout_root = ET.fromstring(source_layout_xml)

            c_sld = new_layout_root.find(_qn(PML_NS, "cSld"))
            new_layout_name = _next_unique_layout_name(base_name, used_layout_name_lowers)
            if c_sld is not None:
                c_sld.attrib["name"] = new_layout_name
            if profile == "title":
                new_layout_root.attrib["type"] = "title"
            elif profile == "content":
                new_layout_root.attrib["type"] = "obj"

            new_layout_xml = ET.tostring(
                new_layout_root, encoding="utf-8", xml_declaration=True
            )
            new_parts[new_layout_part] = new_layout_xml
            new_parts[new_layout_rels_part] = source_layout_rels_xml

            rel_id = next_rel_id()
            ET.SubElement(
                master_rels_root,
                _qn(PKG_REL_NS, "Relationship"),
                {
                    "Id": rel_id,
                    "Type": SLIDE_LAYOUT_REL_TYPE,
                    "Target": f"../slideLayouts/{posixpath.basename(new_layout_part)}",
                },
            )

            while next_layout_id in used_layout_ids:
                next_layout_id += 1
            used_layout_ids.add(next_layout_id)
            ET.SubElement(
                layout_id_list,
                _qn(PML_NS, "sldLayoutId"),
                {
                    "id": str(next_layout_id),
                    _qn(REL_NS, "id"): rel_id,
                },
            )
            next_layout_id += 1

            if new_layout_part not in existing_override_parts:
                ET.SubElement(
                    content_types_root,
                    _qn(CT_NS, "Override"),
                    {
                        "PartName": f"/{new_layout_part}",
                        "ContentType": SLIDE_LAYOUT_CONTENT_TYPE,
                    },
                )
                existing_override_parts.add(new_layout_part)

            _, has_title, has_subtitle, has_body = _layout_features(new_layout_root)
            entries.append(
                {
                    "part": new_layout_part,
                    "rel_id": rel_id,
                    "name_lower": new_layout_name.lower(),
                    "has_title": has_title,
                    "has_subtitle": has_subtitle,
                    "has_body": has_body,
                    "order_index": len(entries),
                }
            )
            created_layout_names.append(new_layout_name)

    if not created_layout_names:
        return {"created_count": 0, "created_layout_names": []}

    updates = {
        master_part: ET.tostring(master_root, encoding="utf-8", xml_declaration=True),
        master_rels_part: ET.tostring(master_rels_root, encoding="utf-8", xml_declaration=True),
        content_types_part: ET.tostring(content_types_root, encoding="utf-8", xml_declaration=True),
    }

    tmp_fd, tmp_raw = tempfile.mkstemp(suffix=".pptx", dir=str(template_path.parent))
    os.close(tmp_fd)
    tmp_path = Path(tmp_raw)
    try:
        with zipfile.ZipFile(template_path, "r") as zin, zipfile.ZipFile(tmp_path, "w") as zout:
            existing_parts: set[str] = set()
            for info in zin.infolist():
                part_name = info.filename
                existing_parts.add(part_name)
                payload = updates.get(part_name)
                if payload is None:
                    payload = zin.read(part_name)
                zout.writestr(info, payload)

            for part_name, payload in new_parts.items():
                if part_name in existing_parts:
                    continue
                zout.writestr(part_name, payload, compress_type=zipfile.ZIP_DEFLATED)
        os.replace(tmp_path, template_path)
    finally:
        if tmp_path.exists():
            tmp_path.unlink()

    return {
        "created_count": len(created_layout_names),
        "created_layout_names": created_layout_names,
    }


def _placeholder_type_name(value: Any) -> str:
    return str(getattr(value, "name", value)).strip().lower()


def _is_content_placeholder(placeholder_type: Any) -> bool:
    if (
        placeholder_type in TITLE_PLACEHOLDERS
        or placeholder_type in SUBTITLE_PLACEHOLDERS
        or placeholder_type in NON_CONTENT_PLACEHOLDERS
    ):
        return False
    return (
        placeholder_type in BODY_PLACEHOLDERS
        or placeholder_type in VISUAL_PLACEHOLDERS
        or placeholder_type in TABLE_PLACEHOLDERS
    )


def _estimate_column_count(placeholder_infos: list[dict[str, Any]], slide_width: int) -> int:
    if not placeholder_infos:
        return 0
    centers = sorted(
        item["left"] + item["width"] // 2
        for item in placeholder_infos
        if item["width"] > 0
    )
    if not centers:
        return 0

    threshold = max(int(slide_width * 0.18), 1)
    columns = [centers[0]]
    for center in centers[1:]:
        if abs(center - columns[-1]) > threshold:
            columns.append(center)
    return len(columns)


def _infer_title_position(
    title_infos: list[dict[str, Any]], slide_width: int, slide_height: int
) -> str:
    if not title_infos:
        return "none"

    title = min(title_infos, key=lambda item: item["top"])
    if title["top"] > int(slide_height * 0.40):
        return "bottom"
    if title["left"] > int(slide_width * 0.55):
        return "right"
    if title["width"] < int(slide_width * 0.40) and title["height"] > int(slide_height * 0.60):
        return "vertical"
    return "top"


def _extract_layout_composition(
    layout: Any,
    *,
    name_lower: str,
    slide_width: int,
    slide_height: int,
) -> dict[str, Any]:
    placeholder_infos: list[dict[str, Any]] = []
    for placeholder in layout.placeholders:
        try:
            placeholder_type = placeholder.placeholder_format.type
        except Exception:
            continue
        placeholder_infos.append(
            {
                "type": placeholder_type,
                "type_name": _placeholder_type_name(placeholder_type),
                "left": int(getattr(placeholder, "left", 0)),
                "top": int(getattr(placeholder, "top", 0)),
                "width": int(getattr(placeholder, "width", 0)),
                "height": int(getattr(placeholder, "height", 0)),
            }
        )

    title_infos = [item for item in placeholder_infos if item["type"] in TITLE_PLACEHOLDERS]
    body_infos = [item for item in placeholder_infos if item["type"] in BODY_PLACEHOLDERS]
    visual_infos = [item for item in placeholder_infos if item["type"] in VISUAL_PLACEHOLDERS]
    table_infos = [item for item in placeholder_infos if item["type"] in TABLE_PLACEHOLDERS]
    content_infos = [item for item in placeholder_infos if _is_content_placeholder(item["type"])]

    column_count = _estimate_column_count(content_infos, slide_width)
    title_position = _infer_title_position(title_infos, slide_width, slide_height)

    caption_like = False
    if visual_infos and body_infos:
        for body in body_infos:
            if body["top"] > int(slide_height * 0.62) or body["height"] < int(slide_height * 0.25):
                caption_like = True
                break

    tags: set[str] = set()
    if "blank" in name_lower:
        tags.add("blank")
    if any(token in name_lower for token in COVER_LAYOUT_KEYWORDS):
        tags.add("cover")
    if "title" in name_lower:
        tags.add("title")
    if any(token in name_lower for token in SECTION_LAYOUT_KEYWORDS):
        tags.add("section")
    if any(token in name_lower for token in VISUAL_LAYOUT_KEYWORDS):
        tags.add("visual_hint")
    if any(token in name_lower for token in TABLE_LAYOUT_KEYWORDS):
        tags.add("table_hint")
    if "caption" in name_lower or caption_like:
        tags.add("caption")
    if "vertical" in name_lower or title_position in {"right", "vertical"}:
        tags.add("vertical")
    if column_count >= 2:
        tags.add("split")
    if len(content_infos) == 1:
        tags.add("single_content")
    if len(content_infos) >= 2:
        tags.add("multi_content")
    if visual_infos:
        tags.add("visual_slot")
    if table_infos:
        tags.add("table_slot")
    if title_infos and not content_infos:
        tags.add("title_only")

    return {
        "placeholder_count": len(placeholder_infos),
        "content_placeholder_count": len(content_infos),
        "body_placeholder_count": len(body_infos),
        "visual_placeholder_count": len(visual_infos),
        "table_placeholder_count": len(table_infos),
        "column_count": column_count,
        "title_position": title_position,
        "tags": sorted(tags),
    }


def list_master_layouts(prs: Presentation) -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    seen_layout_ids: set[int] = set()
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)

    for master_index, master in enumerate(prs.slide_masters):
        for layout_index, layout in enumerate(master.slide_layouts):
            layout_id = id(layout)
            if layout_id in seen_layout_ids:
                continue
            seen_layout_ids.add(layout_id)

            name = str(getattr(layout, "name", "") or "").strip()
            name_lower = name.lower()
            composition = _extract_layout_composition(
                layout,
                name_lower=name_lower,
                slide_width=slide_width,
                slide_height=slide_height,
            )
            placeholder_types: set[Any] = set()
            for placeholder in layout.placeholders:
                try:
                    placeholder_types.add(placeholder.placeholder_format.type)
                except Exception:
                    continue

            items.append(
                {
                    "layout": layout,
                    "name": name,
                    "name_lower": name_lower,
                    "master_index": master_index,
                    "layout_index": layout_index,
                    "has_title": bool(placeholder_types & TITLE_PLACEHOLDERS),
                    "has_subtitle": bool(placeholder_types & SUBTITLE_PLACEHOLDERS),
                    "has_body": bool(placeholder_types & BODY_PLACEHOLDERS),
                    "layout_key": f"{master_index}:{layout_index}:{name_lower}",
                    **composition,
                }
            )

    if not items:
        raise ValueError("Template does not contain slide-master layouts.")
    return items


def build_layout_catalog(layouts: list[dict[str, Any]]) -> list[dict[str, Any]]:
    catalog: list[dict[str, Any]] = []
    for item in layouts:
        catalog.append(
            {
                "master_index": item["master_index"],
                "layout_index": item["layout_index"],
                "layout_name": item["name"],
                "has_title": item["has_title"],
                "has_subtitle": item["has_subtitle"],
                "has_body": item["has_body"],
                "content_placeholder_count": item["content_placeholder_count"],
                "body_placeholder_count": item["body_placeholder_count"],
                "visual_placeholder_count": item["visual_placeholder_count"],
                "table_placeholder_count": item["table_placeholder_count"],
                "column_count": item["column_count"],
                "title_position": item["title_position"],
                "tags": item["tags"],
            }
        )
    return catalog


def inspect_template_layouts(args: argparse.Namespace) -> dict[str, Any]:
    template_path = resolve_template_path(args.template)
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    prs = Presentation(str(template_path))
    layouts = list_master_layouts(prs)
    report = {
        "template_path": str(template_path),
        "layout_count": len(layouts),
        "layouts": build_layout_catalog(layouts),
    }

    if args.layout_report:
        report_path = Path(args.layout_report).expanduser().resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")

    return report


def count_layout_usage(prs: Presentation, layouts: list[dict[str, Any]]) -> dict[str, int]:
    by_layout_id: dict[int, str] = {}
    by_layout_name: dict[str, str] = {}
    for item in layouts:
        key = item["layout_key"]
        by_layout_id[id(item["layout"])] = key
        by_layout_name.setdefault(item["name_lower"], key)

    usage: dict[str, int] = {}
    for slide in prs.slides:
        try:
            slide_layout = slide.slide_layout
        except Exception:
            continue

        layout_key = by_layout_id.get(id(slide_layout))
        if layout_key is None:
            layout_name = str(getattr(slide_layout, "name", "") or "").strip().lower()
            layout_key = by_layout_name.get(layout_name)
        if layout_key is None:
            continue
        usage[layout_key] = usage.get(layout_key, 0) + 1
    return usage


def infer_slide_layout_intent(
    *,
    is_title_slide: bool,
    has_bullets: bool,
    has_table: bool,
    has_visual: bool,
    has_subtitle: bool,
    bullet_count: int,
) -> dict[str, Any]:
    if is_title_slide:
        return {
            "key": "title_cover",
            "description": "Title/cover page with optional subtitle and minimal body placeholders.",
            "target_columns": 0,
            "requires_title": True,
            "requires_body": False,
            "expects_visual": False,
            "expects_table": False,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"cover", "title_only", "title"},
            "discouraged_tags": {"blank", "section", "split", "vertical"},
        }

    if has_table and has_visual:
        return {
            "key": "hybrid_data",
            "description": "Hybrid data page combining visual and tabular information side-by-side.",
            "target_columns": 2,
            "requires_title": True,
            "requires_body": True,
            "expects_visual": True,
            "expects_table": True,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"split", "multi_content", "visual_hint", "table_hint"},
            "discouraged_tags": {"blank", "title_only"},
        }

    if has_table:
        if has_bullets:
            return {
                "key": "table_comparison",
                "description": "Comparison page with table and supporting bullets.",
                "target_columns": 2,
                "requires_title": True,
                "requires_body": True,
                "expects_visual": False,
                "expects_table": True,
                "prefers_subtitle": has_subtitle,
                "preferred_tags": {"table_hint", "split", "multi_content"},
                "discouraged_tags": {"blank", "title_only"},
            }
        return {
            "key": "table_focus",
            "description": "Table-focused page with a single structured data block.",
            "target_columns": 1,
            "requires_title": True,
            "requires_body": True,
            "expects_visual": False,
            "expects_table": True,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"table_hint", "single_content"},
            "discouraged_tags": {"blank", "title_only"},
        }

    if has_visual and has_bullets:
        return {
            "key": "visual_split",
            "description": "Text + visual page with clear split composition.",
            "target_columns": 2,
            "requires_title": True,
            "requires_body": True,
            "expects_visual": True,
            "expects_table": False,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"visual_hint", "split", "multi_content", "caption"},
            "discouraged_tags": {"blank", "title_only"},
        }

    if has_visual:
        return {
            "key": "visual_focus",
            "description": "Visual-led page with optional short caption.",
            "target_columns": 1,
            "requires_title": True,
            "requires_body": True,
            "expects_visual": True,
            "expects_table": False,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"visual_hint", "caption", "visual_slot"},
            "discouraged_tags": {"blank", "split"},
        }

    if has_bullets and bullet_count >= 4:
        return {
            "key": "text_dense",
            "description": "Text-dense page with 4+ bullets requiring generous body area.",
            "target_columns": 1,
            "requires_title": True,
            "requires_body": True,
            "expects_visual": False,
            "expects_table": False,
            "prefers_subtitle": has_subtitle,
            "preferred_tags": {"single_content", "title"},
            "discouraged_tags": {"blank", "title_only", "section"},
        }

    return {
        "key": "text_brief",
        "description": "Short text page with concise bullets.",
        "target_columns": 1,
        "requires_title": True,
        "requires_body": True,
        "expects_visual": False,
        "expects_table": False,
        "prefers_subtitle": has_subtitle,
        "preferred_tags": {"single_content", "section", "title"},
        "discouraged_tags": {"blank", "title_only"},
    }


def _layout_usage_penalty(usage_count: int, best_gap: int) -> int:
    if usage_count <= 0:
        return 0
    if best_gap <= 12:
        return usage_count * 7
    if best_gap <= 22:
        return usage_count * 4
    if best_gap <= 30:
        return usage_count * 2
    return 0


def _score_intent_tags(
    tags: set[str], preferred_tags: set[str], discouraged_tags: set[str]
) -> int:
    score = 0
    for tag in preferred_tags:
        if tag in tags:
            score += 8
    for tag in discouraged_tags:
        if tag in tags:
            score -= 10
    return score


def score_layout_candidate(
    candidate: dict[str, Any],
    *,
    intent: dict[str, Any],
    bullet_count: int,
    has_subtitle: bool,
) -> int:
    name = candidate["name_lower"]
    tags = set(candidate.get("tags", []))
    column_count = int(candidate.get("column_count", 0))
    body_slots = int(candidate.get("body_placeholder_count", 0))
    visual_slots = int(candidate.get("visual_placeholder_count", 0))
    table_slots = int(candidate.get("table_placeholder_count", 0))
    title_position = str(candidate.get("title_position", "none"))
    score = 0

    requires_title = bool(intent.get("requires_title"))
    requires_body = bool(intent.get("requires_body"))
    expects_visual = bool(intent.get("expects_visual"))
    expects_table = bool(intent.get("expects_table"))
    preferred_tags = set(intent.get("preferred_tags", set()))
    discouraged_tags = set(intent.get("discouraged_tags", set()))
    target_columns = int(intent.get("target_columns", 1))
    intent_key = str(intent.get("key", ""))

    if "blank" in tags or "blank" in name:
        score -= 80

    if requires_title and candidate["has_title"]:
        score += 26
    elif requires_title and not candidate["has_title"]:
        score -= 34
    if has_subtitle and candidate["has_subtitle"]:
        score += 16
    elif has_subtitle and not candidate["has_subtitle"]:
        score -= 10

    if requires_body and candidate["has_body"]:
        score += 24
    elif requires_body and not candidate["has_body"]:
        score -= 30

    if expects_visual:
        if visual_slots > 0:
            score += 30
        if "visual_hint" in tags:
            score += 16
        if visual_slots == 0 and "visual_hint" not in tags:
            score -= 24
    else:
        if visual_slots > 0 and intent_key in {"text_dense", "text_brief"}:
            score -= 6

    if expects_table:
        if table_slots > 0:
            score += 30
        if "table_hint" in tags:
            score += 16
        if table_slots == 0 and "table_hint" not in tags:
            score -= 24
    else:
        if table_slots > 0 and intent_key in {"text_dense", "text_brief"}:
            score -= 6

    score += _score_intent_tags(tags, preferred_tags, discouraged_tags)

    if intent_key == "title_cover":
        if candidate["has_body"]:
            score -= 16
        if title_position == "top":
            score += 6
        if "vertical" in tags:
            score -= 8
        return score

    if intent_key == "text_dense" and body_slots >= 2:
        score += 8
    if intent_key == "text_brief" and bullet_count <= 2 and "section" in tags:
        score += 12
    if intent_key == "text_dense" and "section" in tags:
        score -= 10
    if "title and content" in name:
        score += 14

    if target_columns >= 2 and column_count >= 2:
        score += 10
    if target_columns == 1 and column_count == 1:
        score += 6

    if title_position == "bottom" and intent_key in {"text_dense", "text_brief"}:
        score -= 8
    if "vertical" in tags and intent_key in {"text_dense", "text_brief"}:
        score -= 6
    if intent_key == "visual_split" and "split" in tags:
        score += 8
    if intent_key == "visual_focus" and "caption" in tags:
        score += 8
    if intent_key in {"table_comparison", "hybrid_data"} and "split" in tags:
        score += 8
    if requires_body and not candidate["has_body"]:
        score -= 10

    return score


def pick_layout(
    prs: Presentation,
    layouts: list[dict[str, Any]],
    requested_index: Any,
    requested_name: str,
    fallback: int,
    *,
    prefer_master_layouts: bool,
    is_title_slide: bool,
    has_bullets: bool,
    has_table: bool,
    has_visual: bool,
    has_subtitle: bool,
    bullet_count: int = 0,
) -> tuple[Any, str, dict[str, Any]]:
    del requested_index, requested_name, fallback, prefer_master_layouts

    if not layouts:
        raise ValueError("Template does not contain any slide-master layouts.")

    intent = infer_slide_layout_intent(
        is_title_slide=is_title_slide,
        has_bullets=has_bullets,
        has_table=has_table,
        has_visual=has_visual,
        has_subtitle=has_subtitle,
        bullet_count=bullet_count,
    )

    usage_by_layout = count_layout_usage(prs, layouts)
    target_columns = int(intent.get("target_columns", 1))

    scored_base: list[tuple[int, dict[str, Any]]] = []
    for candidate in layouts:
        base_score = score_layout_candidate(
            candidate,
            intent=intent,
            bullet_count=bullet_count,
            has_subtitle=has_subtitle,
        )
        scored_base.append((base_score, candidate))

    best_base_score = max(item[0] for item in scored_base)
    scored: list[tuple[int, int, int, int, int, int, int, int, int, dict[str, Any]]] = []
    for score, candidate in scored_base:
        usage = usage_by_layout.get(candidate["layout_key"], 0)
        usage_penalty = _layout_usage_penalty(usage, best_base_score - score)
        adjusted_score = score - usage_penalty

        # Penalize obvious placeholder mismatches to preserve master intent.
        if is_title_slide and not candidate["has_title"]:
            adjusted_score -= 30
        if has_subtitle and not candidate["has_subtitle"]:
            adjusted_score -= 6
        if (has_bullets or has_table or has_visual) and not candidate["has_body"]:
            adjusted_score -= 12

        column_fit = -abs(int(candidate.get("column_count", 0)) - target_columns)
        scored.append(
            (
                adjusted_score,
                score,
                -usage,
                column_fit,
                int(candidate["has_title"]),
                int(candidate["has_subtitle"]),
                int(candidate["has_body"]),
                -int(candidate["master_index"]),
                -int(candidate["layout_index"]),
                candidate,
            )
        )

    scored.sort(reverse=True)
    selected = scored[0][9]
    return (
        selected["layout"],
        "master_auto",
        {
            "intent_key": intent["key"],
            "intent_description": intent["description"],
            "selected_layout_name": selected["name"],
            "selected_layout_tags": selected["tags"],
        },
    )


def find_body_shape(slide: Any) -> Any | None:
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            placeholder_type = shape.placeholder_format.type
        except Exception:
            continue
        if placeholder_type in BODY_PLACEHOLDERS:
            return shape

    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if getattr(shape, "has_text_frame", False):
            return shape
    return None


def set_title(slide: Any, title_text: str) -> None:
    if not title_text:
        return
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        text_frame = slide.shapes.title.text_frame
        text_frame.clear()
        text_frame.paragraphs[0].text = title_text
        return

    # Keep fallback box safely inside common 10in+ slide widths.
    box = slide.shapes.add_textbox(Inches(0.75), Inches(0.3), Inches(8.5), Inches(0.9))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = title_text
    paragraph.font.size = Pt(30)


def set_subtitle(slide: Any, subtitle_text: str) -> None:
    if not subtitle_text:
        return

    for shape in slide.placeholders:
        if not shape.has_text_frame:
            continue
        try:
            placeholder_type = shape.placeholder_format.type
        except Exception:
            continue
        if placeholder_type == PP_PLACEHOLDER.SUBTITLE:
            text_frame = shape.text_frame
            text_frame.clear()
            text_frame.paragraphs[0].text = subtitle_text
            return

    # Keep fallback subtitle box inside slide bounds.
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.4), Inches(8.5), Inches(0.9))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = subtitle_text
    paragraph.font.size = Pt(20)


def set_bullets(slide: Any, bullets: list[str]) -> None:
    if not bullets:
        return
    body_shape = find_body_shape(slide)
    if body_shape is None:
        body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(6.0), Inches(4.6))
    text_frame = body_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.paragraphs[0].text = bullets[0]
    text_frame.paragraphs[0].level = 0
    for item in bullets[1:]:
        paragraph = text_frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0


def add_table(prs: Presentation, slide: Any, table_spec: dict[str, Any]) -> bool:
    if not table_spec:
        return False

    headers = table_spec.get("headers", [])
    rows = table_spec.get("rows", [])
    if not headers or not rows:
        return False

    column_count = len(headers)
    if column_count == 0:
        return False
    for row in rows:
        if len(row) != column_count:
            raise ValueError("All table rows must match headers length.")

    default_left = int(prs.slide_width * 0.06)
    default_top = int(prs.slide_height * 0.34)
    default_width = int(prs.slide_width * 0.88)
    default_height = int(prs.slide_height * 0.50)

    left = to_emu(table_spec.get("left"), prs.slide_width, default_left)
    top = to_emu(table_spec.get("top"), prs.slide_height, default_top)
    width = to_emu(table_spec.get("width"), prs.slide_width, default_width)
    height = to_emu(table_spec.get("height"), prs.slide_height, default_height)

    shape = slide.shapes.add_table(
        rows=len(rows) + 1,
        cols=column_count,
        left=left,
        top=top,
        width=width,
        height=height,
    )
    table = shape.table

    for col_index, header in enumerate(headers):
        cell = table.cell(0, col_index)
        cell.text = str(header)
        first_paragraph = cell.text_frame.paragraphs[0]
        first_paragraph.font.bold = True

    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = str(value)

    return True


def add_visual(
    prs: Presentation,
    slide: Any,
    visual: dict[str, Any],
    plan_path: Path,
    strict_images: bool,
) -> bool:
    image_ref = visual.get("path")
    if not image_ref:
        return False

    image_path = resolve_path(str(image_ref), plan_path)
    if not image_path.exists():
        message = f"Image not found: {image_path}"
        if strict_images:
            raise FileNotFoundError(message)
        print(f"[WARN] {message}", file=sys.stderr)
        return False

    default_left = int(prs.slide_width * 0.56)
    default_top = int(prs.slide_height * 0.20)
    default_width = int(prs.slide_width * 0.38)

    left = to_emu(visual.get("left"), prs.slide_width, default_left)
    top = to_emu(visual.get("top"), prs.slide_height, default_top)
    width = to_emu(visual.get("width"), prs.slide_width, default_width)
    height_value = visual.get("height")
    height = None if height_value is None else to_emu(height_value, prs.slide_height, 0)

    if height is None:
        picture = slide.shapes.add_picture(str(image_path), left, top, width=width)
    else:
        picture = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    caption = str(visual.get("caption", "")).strip()
    if caption:
        caption_top = picture.top + picture.height + int(Inches(0.05))
        caption_height = int(Inches(0.45))
        box = slide.shapes.add_textbox(left, caption_top, picture.width, caption_height)
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = caption
        paragraph.font.size = Pt(12)

    return True


def add_notes(slide: Any, sources: list[str], speaker_notes: list[str]) -> None:
    clean_sources = [item.strip() for item in sources if item and item.strip()]
    clean_notes = [item.strip() for item in speaker_notes if item and item.strip()]
    if not clean_sources and not clean_notes:
        return

    lines: list[str] = []
    if clean_sources:
        lines.append("Sources:")
        lines.extend(f"- {item}" for item in clean_sources)
    if clean_notes:
        if lines:
            lines.append("")
        lines.append("Notes:")
        lines.extend(f"- {item}" for item in clean_notes)

    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.paragraphs[0].text = "\n".join(lines)


def normalize_slide_spec(raw: Any, index: int) -> dict[str, Any]:
    if not isinstance(raw, dict):
        raise ValueError(f"Slide #{index} must be a JSON object.")

    title = str(raw.get("title", f"Slide {index}")).strip() or f"Slide {index}"
    bullets_raw = raw.get("bullets", [])
    if bullets_raw is None:
        bullets_raw = []
    if not isinstance(bullets_raw, list):
        raise ValueError(f"Slide #{index} field 'bullets' must be a list.")
    bullets = [str(item).strip() for item in bullets_raw if str(item).strip()]

    sources_raw = raw.get("sources", [])
    if sources_raw is None:
        sources_raw = []
    if not isinstance(sources_raw, list):
        raise ValueError(f"Slide #{index} field 'sources' must be a list.")
    sources = [str(item).strip() for item in sources_raw if str(item).strip()]

    notes_raw = raw.get("speaker_notes", [])
    if notes_raw is None:
        notes_raw = []
    if not isinstance(notes_raw, list):
        raise ValueError(f"Slide #{index} field 'speaker_notes' must be a list.")
    speaker_notes = [str(item).strip() for item in notes_raw if str(item).strip()]

    visual_raw = raw.get("visual")
    if visual_raw is None:
        visual = {}
    elif isinstance(visual_raw, dict):
        visual = visual_raw
    else:
        raise ValueError(f"Slide #{index} field 'visual' must be an object.")

    table_raw = raw.get("table")
    if table_raw is None:
        table = {}
    elif isinstance(table_raw, dict):
        headers_raw = table_raw.get("headers", [])
        rows_raw = table_raw.get("rows", [])
        if not isinstance(headers_raw, list):
            raise ValueError(f"Slide #{index} field 'table.headers' must be a list.")
        if not isinstance(rows_raw, list):
            raise ValueError(f"Slide #{index} field 'table.rows' must be a list.")

        headers = [str(item).strip() for item in headers_raw if str(item).strip()]
        rows: list[list[str]] = []
        for row_i, row_item in enumerate(rows_raw, start=1):
            if not isinstance(row_item, list):
                raise ValueError(
                    f"Slide #{index} field 'table.rows[{row_i}]' must be a list."
                )
            row_values = [str(value).strip() for value in row_item]
            if row_values:
                rows.append(row_values)

        if rows and headers:
            expected_cols = len(headers)
            for row_i, row_values in enumerate(rows, start=1):
                if len(row_values) != expected_cols:
                    raise ValueError(
                        f"Slide #{index} field 'table.rows[{row_i}]' must have {expected_cols} columns."
                    )

        table = {
            "headers": headers,
            "rows": rows,
            "left": table_raw.get("left"),
            "top": table_raw.get("top"),
            "width": table_raw.get("width"),
            "height": table_raw.get("height"),
        }
    else:
        raise ValueError(f"Slide #{index} field 'table' must be an object.")

    return {
        "layout": raw.get("layout"),
        "layout_name": str(raw.get("layout_name", "")).strip(),
        "title": title,
        "subtitle": str(raw.get("subtitle", "")).strip(),
        "bullets": bullets,
        "sources": sources,
        "speaker_notes": speaker_notes,
        "visual": visual,
        "table": table,
    }


def add_title_slide_if_requested(
    prs: Presentation,
    layouts: list[dict[str, Any]],
    plan: dict[str, Any],
    fallback_layout: int,
    prefer_master_layouts: bool,
) -> tuple[int, int]:
    title_spec = plan.get("title_slide")
    if not isinstance(title_spec, dict):
        return 0, 0

    requested_layout_name = str(title_spec.get("layout_name", "")).strip()
    layout, selection_mode, _selection_detail = pick_layout(
        prs,
        layouts,
        title_spec.get("layout"),
        requested_layout_name,
        fallback_layout,
        prefer_master_layouts=prefer_master_layouts,
        is_title_slide=True,
        has_bullets=False,
        has_table=False,
        has_visual=False,
        has_subtitle=bool(str(title_spec.get("subtitle", "")).strip()),
    )
    slide = prs.slides.add_slide(layout)
    set_title(slide, str(title_spec.get("title", "")).strip())
    set_subtitle(slide, str(title_spec.get("subtitle", "")).strip())

    title_sources = title_spec.get("sources", [])
    title_notes = title_spec.get("speaker_notes", [])
    if not isinstance(title_sources, list):
        title_sources = []
    if not isinstance(title_notes, list):
        title_notes = []
    add_notes(slide, [str(item) for item in title_sources], [str(item) for item in title_notes])
    master_auto_count = 1 if selection_mode == "master_auto" else 0
    return 1, master_auto_count


def build_deck(args: argparse.Namespace) -> dict[str, Any]:
    template_path = resolve_template_path(args.template)

    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    (
        work_dir,
        staged_template,
        output_path,
        thread_key,
        thread_key_source,
        generated_workdir_key,
        reused_work_dir,
    ) = prepare_work_paths(args, template_path)
    plan, plan_reference_path, staged_plan_path = resolve_plan_input(args.plan, work_dir)

    if not isinstance(plan, dict):
        raise ValueError("Top-level plan JSON must be an object.")
    slides_raw = plan.get("slides")
    if not isinstance(slides_raw, list) or not slides_raw:
        raise ValueError("Plan must contain a non-empty 'slides' list.")
    expected_language = normalize_language(plan.get("language"))
    if plan.get("language") not in (None, "") and expected_language is None:
        raise ValueError("Unsupported plan.language. Use 'ja' or 'en'.")
    resolved_language = infer_language_from_plan(plan, slides_raw)
    auto_slide_order = normalize_auto_slide_order(plan.get("auto_slide_order"))

    for index, raw_slide in enumerate(slides_raw, start=1):
        if not isinstance(raw_slide, dict):
            raise ValueError(f"Slide #{index} must be a JSON object.")

    # The template can be replaced between runs; always inspect and patch per-run.
    need_title_layout = isinstance(plan.get("title_slide"), dict)
    need_content_layout = True
    need_visual_layout = any(
        slide_requests_visual(raw_slide)
        for raw_slide in slides_raw
        if isinstance(raw_slide, dict)
    )
    need_table_layout = any(
        slide_requests_table(raw_slide)
        for raw_slide in slides_raw
        if isinstance(raw_slide, dict)
    )
    layout_bootstrap = ensure_required_master_layouts(
        staged_template,
        need_title_layout=need_title_layout,
        need_content_layout=need_content_layout,
        need_visual_layout=need_visual_layout,
        need_table_layout=need_table_layout,
    )

    # Initialize one output PPTX file, then append slides to it one-by-one.
    if staged_template.resolve() != output_path.resolve() or not output_path.exists():
        shutil.copy2(staged_template, output_path)

    add_slide_script = (Path(__file__).resolve().parent / "add_slide.py").resolve()
    if not add_slide_script.exists():
        raise FileNotFoundError(f"Slide append script not found: {add_slide_script}")

    inserted_visuals = 0
    inserted_tables = 0
    master_auto_layout_count = 0
    added_slides = 0
    layout_intent_summary: dict[str, int] = {}
    layout_selection_trace: list[dict[str, Any]] = []

    def append_one(spec_obj: dict[str, Any], *, kind: str, seq: int) -> dict[str, Any]:
        spec_path = (work_dir / f"slide_spec_{seq:03d}_{kind}.json").resolve()
        spec_path.write_text(json.dumps(spec_obj, ensure_ascii=False, indent=2), encoding="utf-8")

        command = [
            sys.executable,
            str(add_slide_script),
            "--deck",
            str(output_path),
            "--kind",
            kind,
            "--spec",
            str(spec_path),
            "--fallback-layout",
            str(args.fallback_layout),
            "--plan-base",
            str(plan_reference_path.parent),
        ]
        if args.strict_images:
            command.append("--strict-images")
        if args.no_prefer_master_layouts:
            command.append("--no-prefer-master-layouts")
        command.extend(["--language", resolved_language])

        completed = subprocess.run(command, capture_output=True, text=True)
        if completed.returncode != 0:
            stderr_text = (completed.stderr or "").strip()
            stdout_text = (completed.stdout or "").strip()
            detail = stderr_text or stdout_text or "unknown error"
            raise RuntimeError(f"Failed to append {kind} slide #{seq}: {detail}")

        stdout_text = (completed.stdout or "").strip()
        if not stdout_text:
            raise RuntimeError(f"Slide append returned empty output for {kind} slide #{seq}.")
        try:
            return json.loads(stdout_text.splitlines()[-1])
        except json.JSONDecodeError as exc:
            raise RuntimeError(
                f"Slide append returned non-JSON output for {kind} slide #{seq}: {stdout_text}"
            ) from exc

    sequence = 0
    title_spec = plan.get("title_slide")
    sections: dict[str, list[tuple[str, dict[str, Any]]]] = {
        "agenda": [("content", build_agenda_slide(slides_raw, resolved_language))],
        "title": [("title", title_spec)] if isinstance(title_spec, dict) else [],
        "content": [("content", raw_slide) for raw_slide in slides_raw if isinstance(raw_slide, dict)],
        "summary": [("content", build_summary_slide(slides_raw, resolved_language))],
    }

    for section in auto_slide_order:
        for kind, spec_obj in sections.get(section, []):
            slide_result = append_one(spec_obj, kind=kind, seq=sequence)
            sequence += 1
            added_slides += 1
            if slide_result.get("used_master_auto"):
                master_auto_layout_count += 1
            if slide_result.get("table_added"):
                inserted_tables += 1
            if slide_result.get("visual_added"):
                inserted_visuals += 1
            layout_intent = slide_result.get("layout_intent")
            if isinstance(layout_intent, dict):
                intent_key = str(layout_intent.get("intent_key", "")).strip()
                if intent_key:
                    layout_intent_summary[intent_key] = (
                        layout_intent_summary.get(intent_key, 0) + 1
                    )
                layout_selection_trace.append(
                    {
                        "section": section,
                        "kind": kind,
                        "intent_key": intent_key or "unknown",
                        "intent_description": str(
                            layout_intent.get("intent_description", "")
                        ).strip(),
                        "selected_layout_name": str(
                            layout_intent.get("selected_layout_name", "")
                        ).strip(),
                    }
                )

    return {
        "work_dir": work_dir,
        "thread_key": thread_key,
        "thread_key_source": thread_key_source,
        "generated_workdir_key": generated_workdir_key,
        "workdir_key_env": WORKDIR_KEY_ENV,
        "reused_work_dir": reused_work_dir,
        "plan_path": plan_reference_path,
        "staged_plan": staged_plan_path,
        "staged_template": staged_template,
        "output_path": output_path,
        "slide_count": added_slides,
        "visual_count": inserted_visuals,
        "table_count": inserted_tables,
        "master_auto_layout_count": master_auto_layout_count,
        "created_master_layout_count": layout_bootstrap["created_count"],
        "created_master_layouts": layout_bootstrap["created_layout_names"],
        "layout_intent_summary": layout_intent_summary,
        "layout_selection_trace": layout_selection_trace,
        "auto_slide_order": auto_slide_order,
        "language": resolved_language,
    }


def main() -> int:
    args = parse_args()
    try:
        if args.list_layouts:
            report = inspect_template_layouts(args)
            print(json.dumps(report, ensure_ascii=False, indent=2))
            return 0
        result = build_deck(args)
    except Exception as exc:
        print(f"[ERROR] {exc}", file=sys.stderr)
        return 1

    print(f"[OK] Created {result['output_path']}")
    print(f"[OK] Work directory: {result['work_dir']}")
    print(f"[OK] Thread key: {result['thread_key']}")
    print(f"[OK] Thread key source: {result['thread_key_source']}")
    print(f"[OK] Workdir env key: {result['workdir_key_env']}={result['thread_key']}")
    print(f"[OK] Generated new workdir key: {result['generated_workdir_key']}")
    print(f"[OK] Reused work directory: {result['reused_work_dir']}")
    print(f"[OK] Plan source: {result['plan_path']}")
    print(f"[OK] Staged plan: {result['staged_plan']}")
    print(f"[OK] Staged template: {result['staged_template']}")
    print(f"[OK] Slides: {result['slide_count']}")
    print(f"[OK] Slides using auto-selected master layouts: {result['master_auto_layout_count']}")
    print(f"[OK] Slides with inserted tables: {result['table_count']}")
    print(f"[OK] Slides with inserted visuals: {result['visual_count']}")
    print(f"[OK] Auto slide order: {', '.join(result['auto_slide_order'])}")
    print(f"[OK] Created master layouts: {result['created_master_layout_count']}")
    if result["created_master_layouts"]:
        print(f"[OK] Created master layout names: {', '.join(result['created_master_layouts'])}")
    if result["layout_intent_summary"]:
        summary_pairs = [
            f"{key}:{count}"
            for key, count in sorted(result["layout_intent_summary"].items())
        ]
        print(f"[OK] Layout intents: {', '.join(summary_pairs)}")
    print(f"[OK] Language: {result['language']}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
